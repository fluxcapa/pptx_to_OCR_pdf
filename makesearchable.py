# Takes all .pptx's in current directory and makes them into OCR'd searchable PDFs
# Then combines all the PDFs into OneToRuleThemAll.pdf
# Derek Petersen 3/31/2021

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pytesseract import image_to_string
from pytesseract import image_to_pdf_or_hocr
from PIL import Image
from pathlib import Path
from PyPDF2 import PdfFileMerger
import glob
import re
from natsort import natsorted
import shutil

work_path = Path("work")

Path(work_path).mkdir(parents=True, exist_ok=True)

def iter_picture_shapes(pres):
    for slide in pres.slides:
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape

for pres in sorted(glob.iglob(r'*.pptx')):
    print (pres)
    i = 0
    Path("work/"+pres).mkdir(parents=True, exist_ok=True)
    
    # Make separate image files from each pres and store in work dir
    for picture in iter_picture_shapes(Presentation(pres)):
        i += 1
        image = picture.image
        image_bytes = image.blob
        image_filename = 'image{name}.{ext}'.format(name=i, ext=image.ext)
        with open('work/'+pres+'/'+image_filename, 'wb') as f:
            f.write(image_bytes)

    # Make searchable PDFs from images in work dir
    for picture in sorted(glob.iglob(r'work/'+pres+'/*.png')):
        print (picture)
        pdf = image_to_pdf_or_hocr(picture, extension='pdf')
        with open('{name}.pdf'.format(name=picture), 'a+b') as f:
            f.write(pdf)

    # Combine individual PDF pages into original "package" in the current dir
    merger = PdfFileMerger()
    for pdf in natsorted(glob.iglob(r'work/'+pres+'/*.pdf')):
        print (pdf)
        merger.append(pdf)
    merger.write(pres+'.pdf')
    merger.close()

#Combine all completed OCR'd/searchable PDFs into one to rule them all
merger = PdfFileMerger()
for pdf in natsorted(glob.iglob(r'*.pptx.pdf')):
    print (pdf)
    merger.append(pdf)
merger.write('OneToRuleThemAll.pdf')
merger.close()

#clean up work path
shutil.rmtree('work') 
